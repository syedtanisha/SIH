import io
from fastapi import UploadFile, HTTPException, status
from pypdf import PdfReader
import docx
from pptx import Presentation

def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        reader = PdfReader(io.BytesIO(file_bytes))
        text_parts = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_parts.append(t)
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from PDF document: {str(e)}"
        )

def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        doc = docx.Document(io.BytesIO(file_bytes))
        text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                if row_text:
                    text_parts.append(row_text)
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from DOCX document: {str(e)}"
        )

def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())
        return "\n".join(text_parts).strip()
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to extract text from PPTX presentation: {str(e)}"
        )

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="replace").strip()
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Failed to read text file: {str(e)}"
        )

def process_uploaded_document(filename: str, file_bytes: bytes) -> str:
    ext = filename.split(".")[-1].lower() if "." in filename else ""
    if ext == "pdf":
        return extract_text_from_pdf(file_bytes)
    elif ext in ["docx", "doc"]:
        return extract_text_from_docx(file_bytes)
    elif ext in ["pptx", "ppt"]:
        return extract_text_from_pptx(file_bytes)
    elif ext in ["txt", "csv", "md"]:
        return extract_text_from_txt(file_bytes)
    else:
        # Fallback to UTF-8 text decoding
        return extract_text_from_txt(file_bytes)
